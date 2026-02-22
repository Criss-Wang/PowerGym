"""Generate HERON research presentation PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA8, 0x5C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
DARK_CARD = RGBColor(0x25, 0x25, 0x3A)
DARKER_CARD = RGBColor(0x20, 0x20, 0x35)

# ── Layout constants ────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Set corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(3)
        # Bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "  # small triangle
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
    return txBox


def add_card(slide, left, top, width, height, title, body_items,
             title_color=ACCENT_BLUE, card_color=DARK_CARD):
    card = add_shape(slide, left, top, width, height, card_color)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    # Body
    add_bullet_text(slide, left + Inches(0.2), top + Inches(0.55),
                    width - Inches(0.4), height - Inches(0.7),
                    body_items, font_size=13, color=LIGHT_GRAY, bullet_color=title_color)
    return card


def add_circle(slide, left, top, size, fill_color, text="", font_size=12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape


def add_arrow_line(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight connector
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_thin_divider(slide, left, top, width, color=MID_GRAY):
    shape = add_rect(slide, left, top, width, Pt(1), color)
    return shape


# ════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ──────────────────────────────────────────────────────────
# SLIDE 1: TITLE
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Accent bar at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

# Title
add_text_box(slide, MARGIN, Inches(1.8), CONTENT_W, Inches(1.2),
             "HERON", font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, MARGIN, Inches(3.0), CONTENT_W, Inches(0.8),
             "Hierarchical Event-driven Reinforcement-learning Orchestration Network",
             font_size=24, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

add_thin_divider(slide, Inches(4), Inches(3.9), Inches(5.333), ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(4.2), CONTENT_W, Inches(0.6),
             "Agent-Paced Event-Driven Execution for Multi-Agent RL in Cyber-Physical Systems",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, MARGIN, Inches(5.2), CONTENT_W, Inches(0.5),
             "NeurIPS 2026 Datasets & Benchmarks Track",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, MARGIN, Inches(5.7), CONTENT_W, Inches(0.5),
             "Research Team Presentation",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 2: THE PROBLEM
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "The Problem: Simulation-to-Deployment Gap",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(0.6),
             "Existing MARL frameworks assume synchronous, fully-observable environments.\nReal CPS deployments have none of these properties.",
             font_size=18, color=LIGHT_GRAY)

# Three problem cards
card_w = Inches(3.8)
card_h = Inches(3.5)
gap = Inches(0.35)
start_x = MARGIN + Inches(0.15)

# Card 1: Execution mismatch
add_card(slide, start_x, Inches(2.0), card_w, card_h,
         "1. Execution Mismatch",
         [
             "Frameworks use lock-step step()",
             "All agents act simultaneously",
             "Real agents have different tick rates",
             "SCADA: ~2s, field devices: ~100ms",
             "Communication delays ignored",
         ],
         title_color=ACCENT_RED)

# Card 2: Information leakage
add_card(slide, start_x + card_w + gap, Inches(2.0), card_w, card_h,
         "2. Information Leakage",
         [
             "Global state accessible to all agents",
             "No visibility enforcement",
             "Agents trained with info they won't have",
             "\"Global state leak\" in benchmarks",
             "Binary obs: full or nothing",
         ],
         title_color=ACCENT_ORANGE)

# Card 3: Fixed coordination
add_card(slide, start_x + 2 * (card_w + gap), Inches(2.0), card_w, card_h,
         "3. Rigid Coordination",
         [
             "Coordination baked into env code",
             "Can't swap protocols without rewrite",
             "No hierarchical control support",
             "Flat agent models only",
             "No protocol-level experimentation",
         ],
         title_color=ACCENT_PURPLE)

# Bottom tagline
add_text_box(slide, MARGIN, Inches(5.9), CONTENT_W, Inches(0.5),
             "Result: Policies trained in simulation fail in deployment due to unrealistic assumptions",
             font_size=16, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 3: KEY IDEA
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Key Idea",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.0),
             "\"HERON shifts MARL simulation from environment-centric synchronous stepping\n"
             "to agent-paced event-driven execution, making execution model, information\n"
             "structure, and coordination protocols first-class experimental variables.\"",
             font_size=20, color=ACCENT_GREEN, bold=False, alignment=PP_ALIGN.CENTER)

# Three dimension cards
dim_w = Inches(3.8)
dim_h = Inches(3.0)

add_card(slide, start_x, Inches(2.8), dim_w, dim_h,
         "Dimension 1: Execution Model",
         [
             "Sync CTDE mode for RL training",
             "Event-driven mode for testing",
             "Per-agent tick rates & delays",
             "Configurable jitter (Gaussian/Uniform)",
             "CPS-calibrated timing (IEEE 2030)",
         ],
         title_color=ACCENT_BLUE)

add_card(slide, start_x + dim_w + gap, Inches(2.8), dim_w, dim_h,
         "Dimension 2: Information Structure",
         [
             "4-level visibility: public, owner,",
             "  upper_level, system",
             "ProxyAgent mediates ALL state access",
             "Feature-level filtering (not binary)",
             "Ablatable as experimental variable",
         ],
         title_color=ACCENT_GREEN)

add_card(slide, start_x + 2 * (dim_w + gap), Inches(2.8), dim_w, dim_h,
         "Dimension 3: Coordination Protocol",
         [
             "Two-layer: Communication + Action",
             "Vertical: setpoints, price signals",
             "Horizontal: P2P trading, consensus",
             "Swap protocols without agent changes",
             "Compare protocols on same env",
         ],
         title_color=ACCENT_ORANGE)

add_text_box(slide, MARGIN, Inches(6.2), CONTENT_W, Inches(0.5),
             "These three dimensions are orthogonal and independently configurable",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 4: ARCHITECTURE OVERVIEW
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "HERON Architecture",
             font_size=36, color=WHITE, bold=True)

# Agent Hierarchy (left side)
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.4),
             "Agent Hierarchy", font_size=20, color=ACCENT_BLUE, bold=True)

# System Agent box
sys_box = add_shape(slide, Inches(1.0), Inches(1.8), Inches(3.2), Inches(0.7), RGBColor(0x8B, 0x20, 0x20))
add_text_box(slide, Inches(1.1), Inches(1.85), Inches(3.0), Inches(0.6),
             "SystemAgent (L3)\nGlobal Coordination  |  tick: 300s",
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow
add_arrow_line(slide, Inches(2.6), Inches(2.5), Inches(2.6), Inches(2.8), LIGHT_GRAY)

# Coordinator Agent box
coord_box = add_shape(slide, Inches(1.0), Inches(2.8), Inches(3.2), Inches(0.7), RGBColor(0x8B, 0x6B, 0x00))
add_text_box(slide, Inches(1.1), Inches(2.85), Inches(3.0), Inches(0.6),
             "CoordinatorAgent (L2)\nZone Management  |  tick: 60s",
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrows
add_arrow_line(slide, Inches(1.8), Inches(3.5), Inches(1.5), Inches(3.8), LIGHT_GRAY)
add_arrow_line(slide, Inches(3.4), Inches(3.5), Inches(3.7), Inches(3.8), LIGHT_GRAY)

# Field Agents
fa1 = add_shape(slide, Inches(0.5), Inches(3.8), Inches(2.0), Inches(0.7), RGBColor(0x1B, 0x6B, 0x2B))
add_text_box(slide, Inches(0.6), Inches(3.85), Inches(1.8), Inches(0.6),
             "FieldAgent (L1)\ntick: 1s", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

fa2 = add_shape(slide, Inches(2.7), Inches(3.8), Inches(2.0), Inches(0.7), RGBColor(0x1B, 0x6B, 0x2B))
add_text_box(slide, Inches(2.8), Inches(3.85), Inches(1.8), Inches(0.6),
             "FieldAgent (L1)\ntick: 1s", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Core Components (right side)
comp_x = Inches(5.5)
comp_w = Inches(3.5)
comp_h = Inches(0.8)

add_text_box(slide, comp_x, Inches(1.2), Inches(4), Inches(0.4),
             "Core Components", font_size=20, color=ACCENT_GREEN, bold=True)

# ProxyAgent
proxy = add_shape(slide, comp_x, Inches(1.8), comp_w, comp_h, RGBColor(0x1A, 0x4A, 0x7A))
add_text_box(slide, comp_x + Inches(0.15), Inches(1.85), comp_w - Inches(0.3), comp_h,
             "ProxyAgent\nState Cache + Visibility Filter",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# EventScheduler
es = add_shape(slide, comp_x, Inches(2.8), comp_w, comp_h, RGBColor(0x0A, 0x5A, 0x5A))
add_text_box(slide, comp_x + Inches(0.15), Inches(2.85), comp_w - Inches(0.3), comp_h,
             "EventScheduler\nHeap-Based Priority Queue",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# MessageBroker
mb = add_shape(slide, comp_x, Inches(3.8), comp_w, comp_h, RGBColor(0x4A, 0x2A, 0x6A))
add_text_box(slide, comp_x + Inches(0.15), Inches(3.85), comp_w - Inches(0.3), comp_h,
             "MessageBroker\nChannel-Based Routing",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Protocol
pr = add_shape(slide, comp_x, Inches(4.8), comp_w, comp_h, RGBColor(0x6A, 0x3A, 0x0A))
add_text_box(slide, comp_x + Inches(0.15), Inches(4.85), comp_w - Inches(0.3), comp_h,
             "Protocol System\nCommunication + Action Composition",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right side: What's built
add_text_box(slide, Inches(9.5), Inches(1.2), Inches(3.5), Inches(0.4),
             "What's Built", font_size=20, color=ACCENT_ORANGE, bold=True)

built_items = [
    "3-level agent hierarchy",
    "ProxyAgent with 4 visibility scopes",
    "EventScheduler (7 event types)",
    "Dual modes: step() + run_event_driven()",
    "InMemoryBroker (6 message types)",
    "VerticalProtocol + HorizontalProtocol",
    "20 FeatureProviders (power domain)",
    "6 agent types (power domain)",
    "6 test networks (IEEE 13/34/123, ...)",
    "7 tutorial notebooks",
]
add_bullet_text(slide, Inches(9.5), Inches(1.7), Inches(3.5), Inches(5.0),
                built_items, font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN)


# ──────────────────────────────────────────────────────────
# SLIDE 5: PROXYAGENT - STATE MEDIATION
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "ProxyAgent: The State Gatekeeper",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, MARGIN, Inches(1.0), Inches(8), Inches(0.5),
             "ALL state access is mediated through a single point, preventing global state leaks",
             font_size=18, color=LIGHT_GRAY)

# State Cache visualization
cache_x = Inches(0.8)
cache_y = Inches(1.8)
cache_w = Inches(4.5)
cache_h = Inches(4.5)

add_shape(slide, cache_x, cache_y, cache_w, cache_h, DARKER_CARD)
add_text_box(slide, cache_x + Inches(0.2), cache_y + Inches(0.1),
             cache_w, Inches(0.4),
             "State Cache", font_size=18, color=ACCENT_BLUE, bold=True)

# Agent entries in cache
entries = [
    ("battery_1 (L1)", "BatterySOC: 0.503, Capacity: 100.0", ACCENT_GREEN),
    ("battery_2 (L1)", "BatterySOC: 0.498, Capacity: 100.0", ACCENT_GREEN),
    ("coordinator_1 (L2)", "ZoneLoad: 45.2, Aggregated: {...}", ACCENT_ORANGE),
    ("system_agent (L3)", "GridFreq: 60.01, TotalLoad: 120.5", ACCENT_RED),
]

for i, (name, features, color) in enumerate(entries):
    y = cache_y + Inches(0.6) + i * Inches(0.85)
    add_shape(slide, cache_x + Inches(0.15), y, cache_w - Inches(0.3), Inches(0.7), DARK_CARD)
    add_text_box(slide, cache_x + Inches(0.3), y + Inches(0.05),
                 cache_w - Inches(0.6), Inches(0.3),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, cache_x + Inches(0.3), y + Inches(0.35),
                 cache_w - Inches(0.6), Inches(0.3),
                 features, font_size=11, color=MID_GRAY)

# Visibility Filter (middle)
filter_x = Inches(5.8)
filter_y = Inches(2.0)
filter_w = Inches(3.0)
filter_h = Inches(4.0)

add_shape(slide, filter_x, filter_y, filter_w, filter_h, RGBColor(0x3A, 0x1A, 0x1A))
add_text_box(slide, filter_x + Inches(0.15), filter_y + Inches(0.1),
             filter_w, Inches(0.4),
             "Visibility Filter", font_size=18, color=ACCENT_RED, bold=True)
add_text_box(slide, filter_x + Inches(0.15), filter_y + Inches(0.5),
             filter_w - Inches(0.3), Inches(0.35),
             "state.observed_by(requestor_id, level)",
             font_size=12, color=LIGHT_GRAY)

vis_items = [
    ("public", "All agents can see", ACCENT_GREEN),
    ("owner", "Only owning agent", ACCENT_RED),
    ("upper_level", "One level above", ACCENT_ORANGE),
    ("system", "System-level (L3) only", ACCENT_PURPLE),
]
for i, (tag, desc, color) in enumerate(vis_items):
    y = filter_y + Inches(1.1) + i * Inches(0.65)
    add_shape(slide, filter_x + Inches(0.15), y,
              filter_w - Inches(0.3), Inches(0.55), DARK_CARD)
    add_text_box(slide, filter_x + Inches(0.3), y + Inches(0.03),
                 Inches(1.5), Inches(0.25),
                 tag, font_size=13, color=color, bold=True)
    add_text_box(slide, filter_x + Inches(0.3), y + Inches(0.28),
                 filter_w - Inches(0.6), Inches(0.25),
                 desc, font_size=11, color=MID_GRAY)

# Output (right side)
out_x = Inches(9.3)
out_y = Inches(1.8)

add_text_box(slide, out_x, out_y, Inches(3.5), Inches(0.4),
             "Filtered Output", font_size=18, color=ACCENT_GREEN, bold=True)

add_card(slide, out_x, out_y + Inches(0.5), Inches(3.7), Inches(1.8),
         "L1 Agent (battery_1) sees:",
         [
             "Own BatterySOC (public)",
             "Own InternalTemp (owner)",
             "battery_2's BatterySOC (public)",
             "NOT: coordinator's private features",
             "NOT: system-level grid data",
         ],
         title_color=ACCENT_GREEN, card_color=DARK_CARD)

add_card(slide, out_x, out_y + Inches(2.5), Inches(3.7), Inches(1.8),
         "L3 Agent (system) sees:",
         [
             "All public features from all agents",
             "System-level grid frequency",
             "NOT: owner-only features",
             "Full env_context (prices, profiles)",
         ],
         title_color=ACCENT_RED, card_color=DARK_CARD)

# Bottom key insight
add_shape(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.65), RGBColor(0x0A, 0x3A, 0x1A))
add_text_box(slide, MARGIN + Inches(0.2), Inches(6.55), CONTENT_W - Inches(0.4), Inches(0.5),
             "Key Insight: Visibility is an experimental variable, not an implementation detail. "
             "Researchers can ablate over visibility configurations to study information structure impact.",
             font_size=14, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 6: DUAL EXECUTION MODES
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Dual Execution Modes",
             font_size=36, color=WHITE, bold=True)

# Left: CTDE
left_w = Inches(5.8)
add_shape(slide, Inches(0.4), Inches(1.2), left_w, Inches(5.8), DARKER_CARD)
add_text_box(slide, Inches(0.6), Inches(1.3), left_w, Inches(0.4),
             "Mode A: CTDE Training (Synchronous)", font_size=20, color=ACCENT_BLUE, bold=True)

ctde_steps = [
    "1. env.step(actions)",
    "2. Pre-step hook (update load profiles)",
    "3. State sync: all agents reconcile with proxy",
    "4. Phase 1: Hierarchical action application",
    "   SystemAgent \u2192 CoordinatorAgent \u2192 FieldAgent",
    "5. Phase 2: Physics simulation (grid solver)",
    "6. Phase 3: Observation collection (filtered)",
    "7. Phase 4: Reward computation",
    "8. Phase 5: Vectorize + return to RL algorithm",
    "",
    "Properties:",
    "  \u2022 Synchronous: all agents act in one call",
    "  \u2022 Full central access via ProxyAgent",
    "  \u2022 Observations auto-vectorized for RL",
    "  \u2022 Compatible with PettingZoo ParallelEnv",
]
add_bullet_text(slide, Inches(0.6), Inches(1.8), left_w - Inches(0.4), Inches(5.0),
                ctde_steps, font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)

# Right: Event-Driven
right_x = Inches(6.7)
right_w = Inches(6.2)
add_shape(slide, right_x, Inches(1.2), right_w, Inches(5.8), DARKER_CARD)
add_text_box(slide, right_x + Inches(0.2), Inches(1.3), right_w, Inches(0.4),
             "Mode B: Event-Driven Testing (Async)", font_size=20, color=ACCENT_GREEN, bold=True)

ed_steps = [
    "1. env.run_event_driven(analyzer, t_end)",
    "2. EventScheduler processes priority queue:",
    "   AGENT_TICK \u2192 agent.tick()",
    "   MESSAGE_DELIVERY \u2192 async obs/state",
    "   ACTION_EFFECT \u2192 delayed apply_action()",
    "   SIMULATION \u2192 physics update",
    "3. Actions passed via MessageBroker channels",
    "4. Parent-initiated reward cascade",
    "5. Results collected via EventAnalyzer",
    "",
    "Properties:",
    "  \u2022 Asynchronous: per-agent tick rates",
    "  \u2022 Realistic delays (obs, act, msg)",
    "  \u2022 Configurable jitter (Gaussian/Uniform)",
    "  \u2022 CPS-calibrated (IEEE 2030 / NTCIP)",
]
add_bullet_text(slide, right_x + Inches(0.2), Inches(1.8), right_w - Inches(0.4), Inches(5.0),
                ed_steps, font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN)


# ──────────────────────────────────────────────────────────
# SLIDE 7: EVENT-DRIVEN TIMELINE
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Event-Driven Execution Timeline",
             font_size=36, color=WHITE, bold=True)

# Timeline visualization
# Time axis
timeline_y = Inches(1.8)
add_rect(slide, Inches(0.5), timeline_y, Inches(12.3), Pt(2), LIGHT_GRAY)

# Time markers
times = ["t=0", "t=msg_d", "t=2\u00D7msg_d", "t=coord_tick", "t=field_tick",
         "t=field+act_d", "t=wait_int", "t=reward"]
for i, t in enumerate(times):
    x = Inches(0.5) + Inches(i * 1.55) + Inches(0.15)
    add_rect(slide, x + Inches(0.3), timeline_y - Inches(0.1), Pt(2), Inches(0.2), LIGHT_GRAY)
    add_text_box(slide, x - Inches(0.1), timeline_y + Inches(0.1), Inches(1.5), Inches(0.3),
                 t, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Event rows (agent swimlanes)
agents = [
    ("SystemAgent", ACCENT_RED, Inches(2.5)),
    ("ProxyAgent", ACCENT_BLUE, Inches(3.5)),
    ("CoordinatorAgent", ACCENT_ORANGE, Inches(4.5)),
    ("FieldAgent", ACCENT_GREEN, Inches(5.5)),
]

for name, color, y in agents:
    add_text_box(slide, Inches(0.1), y - Inches(0.15), Inches(1.5), Inches(0.3),
                 name, font_size=11, color=color, bold=True)
    add_rect(slide, Inches(0.5), y + Inches(0.15), Inches(12.3), Pt(0.5),
             RGBColor(0x33, 0x33, 0x44))

# Events on the timeline
events_data = [
    # (x_offset, y_agent_idx, label, event_color)
    (0.15, 0, "TICK\nschedule subs\nsend actions", ACCENT_RED),
    (1.7, 1, "MSG\nbuild obs\n+ local_state", ACCENT_BLUE),
    (3.25, 0, "MSG\nrecv obs\ncompute_action", ACCENT_RED),
    (4.8, 2, "TICK\nrecv upstream\ncoordinate\nsend to fields", ACCENT_ORANGE),
    (6.35, 3, "TICK\nrecv upstream\nset_action\nschedule effect", ACCENT_GREEN),
    (7.9, 3, "ACTION_EFFECT\napply_action\nsend state", ACCENT_GREEN),
    (9.45, 0, "SIMULATION\nreq global\nrun physics", ACCENT_RED),
    (11.0, 2, "REWARD\nparent cascade\ncompute reward", ACCENT_ORANGE),
]

for x_off, agent_idx, label, color in events_data:
    x = Inches(0.5 + x_off)
    y = agents[agent_idx][2]
    box = add_shape(slide, x, y - Inches(0.1), Inches(1.3), Inches(0.9), DARK_CARD)
    add_text_box(slide, x + Inches(0.05), y - Inches(0.08),
                 Inches(1.2), Inches(0.85),
                 label, font_size=9, color=color, alignment=PP_ALIGN.CENTER)

# Bottom note
add_shape(slide, MARGIN, Inches(6.6), CONTENT_W, Inches(0.6), RGBColor(0x0A, 0x2A, 0x3A))
add_text_box(slide, MARGIN + Inches(0.2), Inches(6.65), CONTENT_W - Inches(0.4), Inches(0.5),
             "Event priority at same timestamp: ACTION_EFFECT (0) > SIMULATION (1) > MESSAGE_DELIVERY (2)    |    "
             "Per-agent TickConfig: tick_interval, obs_delay, act_delay, msg_delay + jitter",
             font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 8: PROTOCOL SYSTEM
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Composable Protocol System",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(0.5),
             "Two-layer design: separate WHAT to communicate from HOW to coordinate actions",
             font_size=18, color=LIGHT_GRAY)

# Protocol composition formula
add_shape(slide, Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.8), DARK_CARD)
add_text_box(slide, Inches(1.7), Inches(1.8), Inches(10), Inches(0.6),
             "Protocol  =  CommunicationProtocol  +  ActionProtocol\n"
             "               (WHAT to share)               (HOW to coordinate)",
             font_size=18, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Vertical protocols (left)
add_card(slide, Inches(0.5), Inches(2.8), Inches(5.8), Inches(3.8),
         "Vertical Protocols (Top-Down)",
         [
             "VerticalProtocol = NoCommunication + VectorDecomposition",
             "",
             "Parent decomposes joint action vector for subordinates:",
             "  SystemAgent action [a1, a2, a3, a4]",
             "  \u2192 CoordinatorAgent: [a1, a2]",
             "  \u2192 FieldAgent 1: [a3]  |  FieldAgent 2: [a4]",
             "",
             "Domain protocols (planned):",
             "  SetpointProtocol: centralized dispatch",
             "  PriceSignalProtocol: decentralized response",
         ],
         title_color=ACCENT_BLUE)

# Horizontal protocols (right)
add_card(slide, Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.8),
         "Horizontal Protocols (Peer-to-Peer)",
         [
             "HorizontalProtocol = StateShare + NoActionCoordination",
             "",
             "Peers share state and coordinate locally:",
             "  Agent A shares load info with Agent B",
             "  Both adjust independently based on shared info",
             "",
             "Domain protocols (planned):",
             "  P2PTradingProtocol: market clearing",
             "  ConsensusProtocol: gossip-based averaging",
         ],
         title_color=ACCENT_GREEN)

# Key insight
add_shape(slide, MARGIN, Inches(6.8), CONTENT_W, Inches(0.5), RGBColor(0x3A, 0x2A, 0x0A))
add_text_box(slide, MARGIN + Inches(0.2), Inches(6.83), CONTENT_W - Inches(0.4), Inches(0.4),
             "Swap protocols at experiment time \u2014 no agent code changes needed. "
             "Compare centralized dispatch vs. market-based coordination on the same grid.",
             font_size=14, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 9: DATA FLOW SUMMARY
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_PURPLE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Data Flow: Type Transformations",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(0.5),
             "Serialization only at message boundaries \u2014 zero overhead in CTDE training mode",
             font_size=18, color=LIGHT_GRAY)

# Data type cards in a flow
flow_items = [
    ("Agent Layer", "Rich Objects:\nState, Action, Observation\nFeatureProvider instances\nDict[str, FeatureProvider]",
     ACCENT_BLUE, Inches(0.5)),
    ("Proxy Layer", "Object Storage:\nStores State objects directly\nobserved_by() \u2192 filtered\nDict[str, np.ndarray]",
     ACCENT_GREEN, Inches(3.5)),
    ("Message Layer", "Serialized:\nstate.to_dict() / from_dict()\nobs.to_dict() / from_dict()\nAction in payload dict",
     ACCENT_ORANGE, Inches(6.5)),
    ("RL Layer", "Vectors:\nobs.vector() / __array__()\nnp.ndarray for neural nets\nreward: float scalar",
     ACCENT_PURPLE, Inches(9.5)),
]

for title, body, color, x in flow_items:
    add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(2.5), DARK_CARD)
    add_text_box(slide, x + Inches(0.15), Inches(1.9), Inches(2.5), Inches(0.3),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.3), Inches(2.5), Inches(1.8),
                 body, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between
for x in [Inches(3.3), Inches(6.3), Inches(9.3)]:
    add_text_box(slide, x, Inches(2.6), Inches(0.3), Inches(0.4),
                 "\u2192", font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key rules table
table_y = Inches(4.6)
add_text_box(slide, MARGIN, table_y, CONTENT_W, Inches(0.4),
             "Serialization Rules", font_size=18, color=WHITE, bold=True)

rules_left = [
    "Agent \u2192 Proxy (CTDE): direct object reference",
    "Agent \u2192 Proxy (Event): .to_dict() \u2192 message \u2192 .from_dict()",
    "Proxy \u2192 Agent: state.observed_by() \u2192 filtered vectors",
    "Observation \u2192 Policy: __array__() auto-converts",
]
rules_right = [
    "Action \u2192 Broker: Action object in payload dict",
    "Broker \u2192 Agent: msg.payload['action']",
    "Obs response: bundles obs + local_state together",
    "State sync: sync_state_from_observed() reconciles",
]

add_bullet_text(slide, Inches(0.5), table_y + Inches(0.4), Inches(6), Inches(2.5),
                rules_left, font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)
add_bullet_text(slide, Inches(6.8), table_y + Inches(0.4), Inches(6), Inches(2.5),
                rules_right, font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN)


# ──────────────────────────────────────────────────────────
# SLIDE 10: FRAMEWORK COMPARISON
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "HERON vs. Existing Frameworks",
             font_size=36, color=WHITE, bold=True)

# Comparison table
headers = ["Framework", "Abstraction", "Event-Driven", "Visibility", "Protocols", "State\nMediation", "CPS\nFocus"]
data = [
    ["HERON", "Internal env", "Native", "4-level", "Composable", "ProxyAgent", "Yes"],
    ["PettingZoo", "Env-algo API", "No", "Manual", "No", "No", "No"],
    ["EPyMARL", "Algorithm", "No", "No", "No", "No", "No"],
    ["MARLlib", "Algorithm", "No", "No", "No", "No", "No"],
    ["SMAC/v2", "Domain", "No", "Partial", "No", "No", "No"],
    ["Grid2Op", "Power", "No", "N/A", "No", "No", "Yes"],
    ["CityLearn", "Building", "No", "Binary", "No", "No", "Yes"],
]

col_widths = [Inches(1.6), Inches(1.5), Inches(1.4), Inches(1.3), Inches(1.5), Inches(1.5), Inches(1.1)]
table_x = Inches(0.8)
table_y = Inches(1.3)
row_h = Inches(0.55)
header_h = Inches(0.65)

# Header row
x = table_x
for i, (header, w) in enumerate(zip(headers, col_widths)):
    add_shape(slide, x, table_y, w - Inches(0.05), header_h, RGBColor(0x1A, 0x3A, 0x5A))
    add_text_box(slide, x + Inches(0.05), table_y + Inches(0.05),
                 w - Inches(0.15), header_h - Inches(0.1),
                 header, font_size=12, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    x += w

# Data rows
for row_idx, row in enumerate(data):
    x = table_x
    y = table_y + header_h + row_idx * row_h
    bg = DARK_CARD if row_idx % 2 == 0 else DARKER_CARD
    if row_idx == 0:
        bg = RGBColor(0x0A, 0x2A, 0x1A)  # Highlight HERON

    for col_idx, (cell, w) in enumerate(zip(row, col_widths)):
        add_shape(slide, x, y, w - Inches(0.05), row_h - Inches(0.03), bg)
        color = WHITE
        if row_idx == 0:
            color = ACCENT_GREEN
        elif cell in ("No", "N/A"):
            color = RGBColor(0x66, 0x66, 0x77)
        elif cell in ("Yes", "Native", "4-level", "Composable", "ProxyAgent"):
            color = ACCENT_GREEN
        elif cell in ("Partial", "Binary", "Manual"):
            color = ACCENT_ORANGE

        add_text_box(slide, x + Inches(0.05), y + Inches(0.05),
                     w - Inches(0.15), row_h - Inches(0.1),
                     cell, font_size=11, color=color, alignment=PP_ALIGN.CENTER,
                     bold=(row_idx == 0))
        x += w

# Key differentiator
add_shape(slide, MARGIN, Inches(5.8), CONTENT_W, Inches(1.3), DARK_CARD)
add_text_box(slide, MARGIN + Inches(0.3), Inches(5.9), CONTENT_W - Inches(0.6), Inches(0.4),
             "Key Differentiator", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(6.3), CONTENT_W - Inches(0.6), Inches(0.7),
             "PettingZoo standardizes the env\u2194algorithm interface (reset, step, observe, act).\n"
             "HERON standardizes what happens INSIDE the environment (state access, communication, coordination).\n"
             "These are complementary \u2014 HERON envs can export to PettingZoo API.",
             font_size=14, color=LIGHT_GRAY)


# ──────────────────────────────────────────────────────────
# SLIDE 11: CONTRIBUTIONS
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Five Key Contributions",
             font_size=36, color=WHITE, bold=True)

contributions = [
    ("1", "Event-Driven Hierarchical Execution",
     "Dual modes (sync training + event-driven validation) via heap-based EventScheduler.\n"
     "Cannot be achieved by wrapping PettingZoo. Timing is an experimental variable.",
     ACCENT_BLUE),
    ("2", "ProxyAgent for State Mediation",
     "All state access through a single gatekeeper enforcing visibility rules.\n"
     "Prevents the common \"global state leak\" in MARL benchmarks.",
     ACCENT_GREEN),
    ("3", "FeatureProviders with 4-Level Visibility",
     "Granular visibility (public/owner/upper_level/system) as first-class experimental variable.\n"
     "Not binary on/off \u2014 ablate over information structures.",
     ACCENT_ORANGE),
    ("4", "Channel-Isolated MessageBroker",
     "Typed channels with environment isolation for parallel training.\n"
     "Explicit message-based communication matching CPS architectures.",
     ACCENT_PURPLE),
    ("5", "Composable Protocol System",
     "CommunicationProtocol + ActionProtocol composition.\n"
     "Swap coordination mechanisms without changing agent code.",
     ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(contributions):
    y = Inches(1.2) + i * Inches(1.15)
    # Number circle
    add_circle(slide, Inches(0.5), y + Inches(0.1), Inches(0.5), color, num, font_size=18)
    # Title
    add_text_box(slide, Inches(1.2), y, Inches(4.5), Inches(0.35),
                 title, font_size=18, color=color, bold=True)
    # Description
    add_text_box(slide, Inches(1.2), y + Inches(0.35), Inches(11.5), Inches(0.7),
                 desc, font_size=13, color=LIGHT_GRAY)


# ──────────────────────────────────────────────────────────
# SLIDE 12: EXPERIMENT PLAN
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Experiment Plan: 12 Experiments",
             font_size=36, color=WHITE, bold=True)

# Paired experiments
add_text_box(slide, MARGIN, Inches(1.1), CONTENT_W, Inches(0.4),
             "Cross-Domain Paired Experiments (must show consistent ordering)",
             font_size=18, color=ACCENT_ORANGE, bold=True)

paired = [
    ("#1-2", "Visibility Ablation", "system > upper > owner in both domains?", "IEEE 34-bus / 5x5 grid"),
    ("#3-4", "Protocol Comparison", "Vertical vs. horizontal ranking holds?", "Setpoint vs. P2P / Fixed vs. Adaptive"),
    ("#5-6", "Timing Sensitivity", "Sync-to-event gap varies by domain?", "IEEE 2030 / NTCIP 1202 calibrated"),
    ("#7-8", "Algorithm Comparison", "Visibility ordering algorithm-agnostic?", "MAPPO / IPPO / QMIX / TarMAC"),
]

for i, (num, name, question, details) in enumerate(paired):
    y = Inches(1.6) + i * Inches(0.65)
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.55), DARK_CARD)
    add_text_box(slide, Inches(0.6), y + Inches(0.05), Inches(0.7), Inches(0.45),
                 num, font_size=13, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(2.5), Inches(0.45),
                 name, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.8), y + Inches(0.05), Inches(5.0), Inches(0.45),
                 question, font_size=13, color=LIGHT_GRAY)
    add_text_box(slide, Inches(8.8), y + Inches(0.05), Inches(4.0), Inches(0.45),
                 details, font_size=11, color=MID_GRAY)

# Single experiments
add_text_box(slide, MARGIN, Inches(4.4), CONTENT_W, Inches(0.4),
             "Infrastructure Experiments",
             font_size=18, color=ACCENT_BLUE, bold=True)

singles = [
    ("#9", "Scalability", "10 to 2000 agents, log-log plots"),
    ("#10", "Broker Overhead", "Abstraction cost < 5% vs. direct calls"),
    ("#12", "Framework Comparison", "PettingZoo + EPyMARL: LOC and feature gaps"),
]

for i, (num, name, details) in enumerate(singles):
    y = Inches(4.9) + i * Inches(0.55)
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.45), DARKER_CARD)
    add_text_box(slide, Inches(0.6), y + Inches(0.03), Inches(0.7), Inches(0.4),
                 num, font_size=13, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.03), Inches(2.5), Inches(0.4),
                 name, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.8), y + Inches(0.03), Inches(8.5), Inches(0.4),
                 details, font_size=13, color=LIGHT_GRAY)

# Algorithm table
add_text_box(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.3),
             "Algorithms: MAPPO (policy gradient) | IPPO (independent) | QMIX (value decomposition) | TarMAC (learned comm)",
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 13: POWER DOMAIN CASE STUDY
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "Case Study: Power Grid Domain",
             font_size=36, color=WHITE, bold=True)

# Left: Hierarchy
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.4),
             "3-Level Power Grid Hierarchy", font_size=18, color=ACCENT_GREEN, bold=True)

grid_levels = [
    ("GridSystemAgent (L3)", "Grid-level coordination\nManages all microgrids\ntick: 300s", ACCENT_RED, Inches(1.7)),
    ("PowerGridAgent (L2)", "Microgrid coordinator\nManages DERs + loads\ntick: 60s", ACCENT_ORANGE, Inches(2.9)),
    ("DeviceAgent (L1)", "Battery, Generator, ESS\nLocal control + actuation\ntick: 1s", ACCENT_GREEN, Inches(4.1)),
]

for name, desc, color, y in grid_levels:
    add_shape(slide, Inches(0.5), y, Inches(4.8), Inches(1.0), DARK_CARD)
    add_text_box(slide, Inches(0.7), y + Inches(0.05), Inches(4.4), Inches(0.3),
                 name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(0.7), y + Inches(0.35), Inches(4.4), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY)

# Right: Feature summary
add_text_box(slide, Inches(5.8), Inches(1.2), Inches(7), Inches(0.4),
             "20 FeatureProviders", font_size=18, color=ACCENT_BLUE, bold=True)

feature_categories = [
    ("Device Features", "BatterySOC, GeneratorOutput, ESSCharge,\nSolarIrradiance, WindSpeed, TransformerLoad",
     ACCENT_GREEN),
    ("Grid Features", "BusVoltage, LineFlow, LineLoading,\nFrequency, PowerBalance, LossMetric",
     ACCENT_ORANGE),
    ("Market Features", "ElectricityPrice, DemandResponse,\nGridTariff, P2PTradeStatus",
     ACCENT_PURPLE),
    ("Environmental", "WeatherCondition, LoadProfile,\nTimeOfDay, SeasonalPattern",
     ACCENT_BLUE),
]

for i, (cat, features, color) in enumerate(feature_categories):
    y = Inches(1.7) + i * Inches(1.0)
    add_shape(slide, Inches(5.8), y, Inches(7.0), Inches(0.85), DARK_CARD)
    add_text_box(slide, Inches(6.0), y + Inches(0.05), Inches(2.2), Inches(0.3),
                 cat, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(8.2), y + Inches(0.05), Inches(4.4), Inches(0.75),
                 features, font_size=11, color=LIGHT_GRAY)

# Networks
add_text_box(slide, Inches(5.8), Inches(5.7), Inches(7), Inches(0.3),
             "6 Test Networks", font_size=18, color=ACCENT_ORANGE, bold=True)
networks = ["IEEE 13-bus", "IEEE 34-bus", "IEEE 123-bus", "CIGRE LV", "Case34 3-phase", "Case LVMG"]
add_text_box(slide, Inches(5.8), Inches(6.1), Inches(7), Inches(0.4),
             "  |  ".join(networks), font_size=13, color=LIGHT_GRAY)


# ──────────────────────────────────────────────────────────
# SLIDE 14: WHAT MAKES THIS NOVEL
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_PURPLE)

add_text_box(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
             "What Makes HERON Novel?",
             font_size=36, color=WHITE, bold=True)

# Three key arguments
novelty_items = [
    ("Different Abstraction Level",
     "PettingZoo standardizes the env\u2194algorithm INTERFACE.\n"
     "HERON standardizes what happens INSIDE the environment.\n"
     "Event-driven execution cannot be achieved by wrapping \u2014\n"
     "it requires changing the fundamental step loop.",
     ACCENT_BLUE, Inches(1.2)),
    ("Three Orthogonal Experimental Dimensions",
     "1. Execution model: sync vs. event-driven (with CPS-calibrated timing)\n"
     "2. Information structure: 4-level feature visibility (ablatable)\n"
     "3. Coordination protocol: composable Communication + Action layers\n"
     "Each is independently configurable. No other framework offers this.",
     ACCENT_GREEN, Inches(3.2)),
    ("Explicit Non-Claims (Intellectual Honesty)",
     "NOT \"privacy-preserving\" in the DP/cryptographic sense.\n"
     "We don't claim to \"discover\" findings \u2014 we provide infrastructure.\n"
     "We COMPLEMENT PettingZoo (different level), not replace it.\n"
     "HERON envs can export to PettingZoo's ParallelEnv API.",
     ACCENT_ORANGE, Inches(5.2)),
]

for title, body, color, y in novelty_items:
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.7), DARK_CARD)
    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(11.7), Inches(0.35),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.5), Inches(11.7), Inches(1.1),
                 body, font_size=14, color=LIGHT_GRAY)


# ──────────────────────────────────────────────────────────
# SLIDE 15: THANK YOU / DISCUSSION
# ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(1.0),
             "Discussion", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_thin_divider(slide, Inches(4.5), Inches(3.2), Inches(4.333), ACCENT_BLUE)

discussion_items = [
    "How should we prioritize the traffic domain vs. deeper power experiments?",
    "Which 4 algorithms should we commit to? (MAPPO, IPPO, QMIX, TarMAC)",
    "Stretch goal: Game domain appendix \u2014 is it worth the effort?",
    "Timeline: 8 weeks to NeurIPS deadline \u2014 are we on track?",
]
add_bullet_text(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(3.0),
                discussion_items, font_size=18, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)

add_text_box(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.5),
             "NeurIPS 2026 Datasets & Benchmarks Track",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "HERON_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
